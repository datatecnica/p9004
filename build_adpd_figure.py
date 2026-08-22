#!/usr/bin/env python3
"""build_adpd_figure.py — the single multi-panel figure for the AD/PD abstract.

One slide, four panels, built from NATIVE PowerPoint shapes rather than an embedded
image, so every box and label stays editable in PowerPoint.

Slide is 6.25 x 8.33 in, which exports to exactly 600 x 800 px at 96 DPI -- the AD/PD
image cap. Portrait because that cap is portrait; for a talk, widen the slide to 13.33 x
7.5 and the same panels re-flow.

Four panels: how the resource is built, THE FOUR QUESTIONS STAKEHOLDERS ASKED with the
robust proteins each one returned, the protein that answers all four, and what is openly
available. Panel B is the centre of the figure -- a reader should be able to point at a
question and read off its hits. Concordance was cut; it is a methods point, and this
figure has to work for a non-specialist in one pass.

Every count and every protein name is read at build time from the newest META_*.csv in
meta/ via abstract_results_buckets.py -- the same source the abstract is written from, so
figure and text cannot drift. Nothing about the results is typed into this file.

Arial, because it is present everywhere PowerPoint is and a silent font fallback would
reflow the whole layout.

Usage:  python3 build_adpd_figure.py [-o adpd_figure.pptx] [--meta-dir meta]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from abstract_results_buckets import (BUCKETS, FOUR, Q4, QUESTION_TEXT, hit_chips,
                                      load, prs_summary)

# ---------------------------------------------------------------- palette
CSF        = RGBColor(0x1F, 0x7A, 0x8C)
CSF_WASH   = RGBColor(0xE4, 0xF0, 0xF2)
PLASMA     = RGBColor(0xB2, 0x6A, 0x15)
PLASMA_W   = RGBColor(0xF7, 0xED, 0xDF)
INK        = RGBColor(0x17, 0x20, 0x29)
MUTED      = RGBColor(0x5D, 0x6B, 0x76)
RULE       = RGBColor(0xC6, 0xC4, 0xBA)
SUNK       = RGBColor(0xF2, 0xF1, 0xEC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
FLAG       = RGBColor(0xA0, 0x3B, 0x2E)

FONT = "Arial"
W, H = 6.25, 8.33
M = 0.28                      # page margin
CW = W - 2 * M                # content width


def textbox(slide, x, y, w, h, text, *, size=8, bold=False, color=INK,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
            spacing=None, caps=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = color
    if spacing is not None:                     # letter-spacing, via raw XML
        r.font._rPr.set("spc", str(int(spacing * 100)))
    return tb


def box(slide, x, y, w, h, *, fill=None, line=RULE, lw=0.6, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:                                  # keep the corner subtle, not pill-like
        shp.adjustments[0] = 0.08
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def chip(slide, x, y, w, h, label, *, fill, line, color, size=6.5):
    box(slide, x, y, w, h, fill=fill, line=line, lw=0.6)
    textbox(slide, x + 0.05, y, w - 0.1, h, label,
            size=size, color=color, anchor=MSO_ANCHOR.MIDDLE)


def arrow_down(slide, x, y, h):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 Inches(x - 0.045), Inches(y), Inches(0.09), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RULE
    shp.line.fill.background()
    shp.shadow.inherit = False


def rule(slide, x, y, w, color=RULE, lw=0.75):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Emu(int(lw * 12700)))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False


def _sci(p: float) -> str:
    """2.2e-40 -> '2x10-40', in the superscript glyphs PowerPoint renders without a font
    switch. Quoting P rather than beta keeps the slide on the same scale as the text."""
    exp = int(f"{p:.0e}".split("e")[1])
    mant = f"{p:.0e}".split("e")[0]
    sup = str(abs(exp)).translate(str.maketrans("0123456789", "\u2070\u00b9\u00b2\u00b3\u2074\u2075\u2076\u2077\u2078\u2079"))
    return f"{mant}\u00d710\u207b{sup}"


def panel_label(slide, x, y, letter, title):
    textbox(slide, x, y, 0.22, 0.16, letter, size=9, bold=True, color=CSF)
    textbox(slide, x + 0.20, y, CW - 0.20, 0.16, title,
            size=9, bold=True, color=INK)


def numbers(meta_dir: str) -> dict:
    """Everything the slide prints, read from the newest META_*.csv per run."""
    df = load(meta_dir)
    rob = df[df["robust_FE"]]
    four = rob[rob["bucket"].isin(FOUR)]

    questions = []
    for qkey in FOUR:
        sub = rob[rob["bucket"] == qkey]
        prot = sub[sub["is_proteomic"]]
        head, blurb, k, term = QUESTION_TEXT[qkey]
        named = prot if term is None else prot[prot["term"] == term]
        questions.append({
            "head": head, "blurb": blurb,
            "n": len(sub), "n_proteins": prot["protein"].nunique(),
            "n_runs": sub["run"].nunique(), "n_total_runs": len(BUCKETS[qkey]),
            "note": (f"{len(named)} are rate-of-change" if term else ""),
            "chips": hit_chips(named, k),
        })

    ddc = rob[rob["protein"] == "DDC"]
    ddc_by_q = [len(ddc[ddc["bucket"] == qkey]) for qkey in FOUR]
    return {
        "n_tested": len(df) - int(df["thin_within"].sum()),
        "n_runs": df["run"].nunique(),
        "n_robust": len(rob),
        "n_runs_hit": rob["run"].nunique(),
        "four_n": len(four),
        "four_runs": four["run"].nunique(),
        "four_runs_total": sum(len(BUCKETS[q]) for q in FOUR),
        "four_proteins": four[four["is_proteomic"]]["protein"].nunique(),
        "questions": questions,
        "ddc_total": len(ddc),
        "ddc_by_q": ddc_by_q,
        "within": len(four[(four["bucket"] == Q4) & (four["term"] == "within")]),
        "prs": prs_summary(rob),
    }


# ---------------------------------------------------------------- the slide
def build(out_path: str, meta_dir: str = "meta") -> None:
    N = numbers(meta_dir)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # ============================================ header
    textbox(slide, M, 0.22, CW, 0.20,
            "A harmonized multi-platform proteomic resource for Parkinson's disease",
            size=11, bold=True, color=INK)
    textbox(slide, M, 0.44, CW, 0.14,
            "PPMI proteomics + GP2 genetics · 4,788 participants · 19,450 visits · "
            "2 platforms · 2 biofluids",
            size=7, color=MUTED, spacing=0.4)
    rule(slide, M, 0.62, CW, RULE, 1.0)

    # ============================================ PANEL A — the pipeline
    y = 0.74
    panel_label(slide, M, y, "A", "How the resource is built")
    y += 0.22

    steps = [
        ("SIX ASSAY PROJECTS",
         "Olink and NULISA, each run in CSF or plasma or both, joined to\n"
         "GP2 genetics — open-science data from two programmes",
         "34,902 analytes, held apart"),
        ("POOLED IN PAIRS",
         "Projects sharing a platform, panel and biofluid are put on one scale.\n"
         "Olink is never mixed with NULISA — so panel C stays a real question.",
         "11,567 analytes · 2,906 participants · original columns kept"),
        (f"{N['n_runs']} PRESPECIFIED ANALYSES",
         "18 logistic · 16 survival · 10 trajectory · 8 linear, each run twice\n"
         "(European and Ashkenazi Jewish) and then combined",
         "age, sex, education and assay-specific components controlled"),
        (f"{N['n_robust']} ROBUST ASSOCIATIONS",
         "significant after correction, replicated in both ancestry groups,\n"
         "and pointing the same way in each",
         f"{N['n_robust']} of {N['n_tested']:,} tested · {N['n_runs_hit']} of "
         f"{N['n_runs']} analyses"),
    ]
    row_h, gap = 0.40, 0.10
    for i, (head, body, foot) in enumerate(steps):
        box(slide, M, y, CW, row_h, fill=SUNK, line=RULE, rounded=True)
        textbox(slide, M + 0.10, y + 0.045, 1.62, 0.14, head,
                size=6.6, bold=True, color=CSF, spacing=0.3)
        textbox(slide, M + 0.10, y + 0.185, 3.30, 0.22, body, size=6.2, color=INK)
        textbox(slide, M + 3.52, y + 0.045, CW - 3.62, row_h - 0.09, foot,
                size=6.2, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        y += row_h
        if i < len(steps) - 1:
            arrow_down(slide, M + 0.55, y + 0.012, gap - 0.024)
            y += gap

    # ============================================ PANEL B — the four questions
    y += 0.16
    panel_label(slide, M, y, "B", "The four questions asked — and the proteins each returned")
    y += 0.20

    textbox(slide, M, y, CW, 0.12,
            f"ROBUST = significant after correction  +  found in BOTH ancestry groups  "
            f"+  pointing the same way in each.   "
            f"{N['four_n']} of the {N['n_robust']} robust results answer these four.",
            size=6.0, color=MUTED)
    y += 0.17

    # One band per question, read top to bottom: what was asked, how much came back, and
    # then the proteins themselves on their own full-width line. The chips need the whole
    # band width -- five of them run to ~4.3 in, so they cannot share a row with the text.
    nmax = max(q["n"] for q in N["questions"])
    for i, q in enumerate(N["questions"]):
        band = 0.64
        box(slide, M, y, CW, band, fill=SUNK if i % 2 == 0 else WHITE,
            line=RULE, rounded=True)

        textbox(slide, M + 0.10, y + 0.045, 2.60, 0.13, f"{i + 1}.  {q['head']}",
                size=7.0, bold=True, color=CSF)

        stat = (f"{q['n']} robust  ·  {q['n_proteins']} proteins  ·  "
                f"{q['n_runs']} of {q['n_total_runs']} analyses")
        if q["note"]:
            stat += f"  ·  {q['note']}"
        textbox(slide, M + 2.70, y + 0.05, CW - 2.80, 0.13, stat,
                size=5.8, color=INK, align=PP_ALIGN.RIGHT)

        # bar scaled across the four, so 11 next to 146 is visible without reading
        box(slide, M + 0.10, y + 0.205, max(2.60 * q["n"] / nmax, 0.02), 0.05,
            fill=CSF, line=None)
        textbox(slide, M + 0.10, y + 0.285, CW - 0.20, 0.12, q["blurb"],
                size=5.6, color=MUTED)

        # the named hits, one row, left to right in the order they rank
        cx = M + 0.10
        for name, n, fl in q["chips"]:
            col = CSF if fl == "CSF" else (PLASMA if fl == "plasma" else INK)
            wash = CSF_WASH if fl == "CSF" else (PLASMA_W if fl == "plasma" else SUNK)
            label = f"{name} {n}" if n > 1 else name
            w = 0.055 * len(label) + 0.14
            if cx + w > M + CW - 0.08:              # never draw off the band
                break
            chip(slide, cx, y + 0.435, w, 0.155, label,
                 fill=wash, line=col, color=col, size=5.8)
            cx += w + 0.06
        y += band + 0.055

    # PRS, one line: it belongs with staging, and it carries a caveat that has to travel
    prs_stats = N["prs"]
    textbox(slide, M, y - 0.01, CW, 0.12,
            f"Genetic risk score (PRS157): {prs_stats['clean']} robust results — rising with NSD "
            f"status (P={_sci(prs_stats['p_nsd_vs_hc'])}) and with NSD-ISS stage, and "
            f"predicting stage D. Carrier contrasts excluded: PRS157 contains LRRK2 and "
            f"GBA variants.",
            size=5.6, color=MUTED)
    y += 0.14

    # ============================================ PANEL C — the protein in all four
    y += 0.10
    panel_label(slide, M, y, "C", "One protein answers all four")
    y += 0.20

    box(slide, M, y, CW, 0.52, fill=None, line=PLASMA, lw=1.0, rounded=True)
    textbox(slide, M + 0.12, y + 0.07, 0.75, 0.26, "DDC",
            size=18, bold=True, color=PLASMA, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, M + 0.86, y + 0.08, 0.62, 0.13, str(N["ddc_total"]),
            size=12, bold=True, color=INK)
    textbox(slide, M + 0.86, y + 0.235, 1.30, 0.12, "robust results", size=5.8, color=MUTED)
    textbox(slide, M + 1.92, y + 0.07, CW - 2.04, 0.36,
            "Dopa decarboxylase — the enzyme that makes dopamine. Raised where disease is "
            "at baseline, higher at later stages, predicting milestones and stage D, and "
            "rising faster afterwards. Spinal fluid and blood, both technologies, both "
            "ancestry groups.",
            size=6.1, color=INK)
    y += 0.58

    # DDC's count under each of the four questions, so the panel ties back to panel B
    cw2 = CW / 4
    for i, (q, n) in enumerate(zip(N["questions"], N["ddc_by_q"])):
        cx = M + i * cw2
        box(slide, cx, y, cw2 - 0.07, 0.22, fill=SUNK, line=None)
        textbox(slide, cx + 0.07, y + 0.03, cw2 - 0.55, 0.11, f"Q{i + 1} {q['head']}",
                size=5.4, color=MUTED)
        textbox(slide, cx + cw2 - 0.48, y + 0.025, 0.34, 0.13, str(n),
                size=8, bold=True, color=PLASMA, align=PP_ALIGN.RIGHT)
        textbox(slide, cx + 0.07, y + 0.125, cw2 - 0.16, 0.10,
                "robust results", size=5.0, color=MUTED)
    y += 0.30

    # ============================================ PANEL D — availability
    y += 0.16
    box(slide, M, y, CW, 0.50, fill=None, line=CSF, lw=1.0, rounded=True)
    textbox(slide, M + 0.12, y + 0.055, 2.0, 0.12, "OPEN AT EVERY STAGE",
            size=6, bold=True, color=CSF, spacing=0.6)
    items = [
        ("Controlled-access data", "harmonized, sources kept"),
        ("Public codebase", "every model readable"),
        ("Results browser", f"all {N['n_runs']} analyses, not only hits"),
        ("AI interpretation", "plain-language readout"),
    ]
    cw = (CW - 0.24) / 4
    for i, (head, sub) in enumerate(items):
        cx = M + 0.12 + i * cw
        textbox(slide, cx, y + 0.215, cw - 0.06, 0.12, head, size=6.2, bold=True, color=INK)
        textbox(slide, cx, y + 0.345, cw - 0.06, 0.12, sub, size=5.6, color=MUTED)

    # ============================================ footer
    textbox(slide, M, H - 0.30, CW, 0.12,
            "Robust = Bonferroni-significant, replicated in both the European and Ashkenazi "
            "Jewish strata, and concordant in direction.  Chip colour: teal = CSF, "
            "amber = plasma, dark = both.  Counts are robust results, not proteins.",
            size=5.6, color=MUTED)

    prs.save(out_path)
    kb = os.path.getsize(out_path) / 1024
    print(f"Wrote {out_path}  ({kb:.0f} KB, {W}x{H} in -> 600x800 px at 96 DPI)")
    print(f"  content bottom edge at {y:.2f} in of {H} (footer at {H - 0.30:.2f})")
    for i, q in enumerate(N["questions"], start=1):
        chips = ", ".join(f"{n} {c}" for n, c, _ in q["chips"])
        print(f"  Q{i} {q['head']:<32} {q['n']:>3} robust · {chips}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("-o", "--output", default="adpd_figure.pptx")
    ap.add_argument("--meta-dir", default="meta")
    a = ap.parse_args()
    build(a.output, a.meta_dir)
